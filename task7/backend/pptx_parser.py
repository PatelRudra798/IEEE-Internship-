# pptx_parser.py
"""Utility to extract plain text from PowerPoint (.pptx) files using python-pptx."""

from pathlib import Path
from pptx import Presentation

def extract_text_from_pptx(file_path: str) -> str:
    """Return concatenated slide text from a PPTX file.

    Args:
        file_path: Path to the PPTX file.
    Returns:
        Plain text extracted from all slide shapes.
    """
    path = Path(file_path)
    if not path.is_file():
        raise FileNotFoundError(f"PPTX not found: {file_path}")
    prs = Presentation(str(path))
    texts = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_texts.append(shape.text)
        texts.append("\n".join(slide_texts))
    return "\n\n---\n\n".join(texts)
